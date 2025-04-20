from fastapi import APIRouter, UploadFile, File, HTTPException, Body
from typing import Dict, Tuple
import PyPDF2
import io
import pytesseract
from PIL import Image
from pptx import Presentation
import os
import tempfile
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
import shutil
import logging
from dotenv import load_dotenv
from pathlib import Path
from typing import Any, List, Optional, Mapping
import cloudinary
import cloudinary.uploader
import uuid
import re # Import regex

# Configure logging first
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Load environment variables from backend/.env
env_path = Path(__file__).parent.parent / '.env'
load_dotenv(dotenv_path=env_path)

# Configure Cloudinary
try:
    cloudinary.config(
        cloud_name=os.getenv("CLOUDINARY_CLOUD_NAME"),
        api_key=os.getenv("CLOUDINARY_API_KEY"),
        api_secret=os.getenv("CLOUDINARY_API_SECRET")
    )
    logger.info("Cloudinary configured successfully.")
except Exception as e:
    logger.error(f"Failed to configure Cloudinary: {str(e)}")
    # Depending on the desired behavior, you might want to raise an error or handle it differently.
    # For now, we log the error and continue, endpoints requiring Cloudinary might fail.

router = APIRouter(
    prefix="/api/rag",
    tags=["RAG"],
    responses={404: {"description": "Not found"}},
)

# Initialize text splitter for chunking
text_splitter = RecursiveCharacterTextSplitter(
    chunk_size=1000,
    chunk_overlap=200,
    length_function=len,
)

# Initialize Google Gemini LLM
try:
    llm = ChatGoogleGenerativeAI(
        model="gemini-2.0-flash-lite",
        google_api_key=os.getenv("GOOGLE_API_KEY"),
        temperature=0.7,
        convert_system_message_to_human=True
    )
    logger.info("Successfully initialized Google Gemini LLM")
except Exception as e:
    logger.error(f"Failed to initialize Google Gemini LLM: {str(e)}")
    raise

# Initialize embeddings using a smaller, non-gated model
cache_dir = os.path.join(os.path.dirname(__file__), '..', 'model_cache') # Adjusted cache path
os.makedirs(cache_dir, exist_ok=True)

# Replace the existing embeddings code with:
embeddings = HuggingFaceEmbeddings(
    model_name="sentence-transformers/all-MiniLM-L6-v2",  # Public model, no auth required
    model_kwargs={'device': 'cpu'},
    encode_kwargs={
        'normalize_embeddings': True,
        'batch_size': 32
    },
    cache_folder=cache_dir
)

# --- Updated Vector Store Processing ---
async def process_text_to_vectors(text: str, user_id: str, document_id: str) -> Dict:
    """Process text into chunks and create/update vector store for a specific document"""
    try:
        chunks = text_splitter.split_text(text)
        if not chunks:
             logger.warning(f"No text chunks generated for user_id: {user_id}, document_id: {document_id}")
             return {"status": "success", "message": "No text content found to process."}

        # Document-specific vector store path
        vector_store_dir = f"vector_stores/{user_id}/{document_id}"
        vector_store_path = os.path.join(vector_store_dir, "vectors")
        os.makedirs(vector_store_dir, exist_ok=True)

        # Create a new vector store for this document (overwrite if exists)
        logger.info(f"Creating new vector store for user_id: {user_id}, document_id: {document_id}")
        vectorstore = FAISS.from_texts(chunks, embeddings)
        logger.info(f"Created vector store with {len(chunks)} chunks for document_id: {document_id}")

        vectorstore.save_local(vector_store_path)
        logger.info(f"Vector store saved successfully for document_id: {document_id} at {vector_store_path}")

        return {"status": "success", "chunks_processed": len(chunks)}
    except Exception as e:
        logger.error(f"Error processing text for user_id {user_id}, document_id {document_id}: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error processing text: {str(e)}")

# --- Updated Cloudinary Upload ---
async def upload_to_cloudinary(file_content: bytes, user_id: str, filename: str) -> Tuple[str, str]:
    """Uploads file content to Cloudinary and returns secure_url and public_id"""
    try:
        # Generate a unique identifier part for the public_id
        unique_id = str(uuid.uuid4())
        # Sanitize filename for public_id (remove potentially problematic characters)
        sanitized_filename = re.sub(r'[^a-zA-Z0-9_.-]', '_', filename)
        public_id = f"notes/{user_id}/{unique_id}_{sanitized_filename}"

        logger.info(f"Uploading file to Cloudinary with public_id: {public_id}")
        upload_result = cloudinary.uploader.upload(
            file_content,
            public_id=public_id,
            resource_type="auto"
        )
        secure_url = upload_result['secure_url']
        returned_public_id = upload_result['public_id'] # Use the public_id returned by Cloudinary
        logger.info(f"File uploaded successfully to Cloudinary: {secure_url}, Public ID: {returned_public_id}")
        return secure_url, returned_public_id
    except Exception as e:
        logger.error(f"Cloudinary upload failed for user {user_id}, filename {filename}: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Cloudinary upload error: {str(e)}")

def extract_document_id_from_public_id(public_id: str) -> str:
    """Extracts a unique document identifier from the Cloudinary public_id"""
    # Assumes public_id format like "notes/user_id/uuid_filename"
    parts = public_id.split('/')
    if len(parts) > 2:
        # Extract the part after the user_id, e.g., "uuid_filename"
        # Or even better, just the UUID if it's always present at the start
        filename_part = parts[-1]
        uuid_match = re.match(r'^[0-9a-fA-F]{8}-[0-9a-fA-F]{4}-[0-9a-fA-F]{4}-[0-9a-fA-F]{4}-[0-9a-fA-F]{12}', filename_part)
        if uuid_match:
            return uuid_match.group(0) # Return just the UUID
        else:
            # Fallback to using the last part if UUID is not found (less ideal)
            logger.warning(f"Could not extract UUID from public_id: {public_id}. Using last part.")
            return filename_part
    logger.error(f"Could not extract document_id from unexpected public_id format: {public_id}")
    # Return a default or raise an error if extraction fails critically
    return public_id.replace('/', '_') # Basic fallback


@router.post("/process_pdf")
async def process_pdf(user_id: str = Body(...), file: UploadFile = File(...)):
    if not file.filename.endswith('.pdf'):
        raise HTTPException(status_code=400, detail="File must be a PDF")
    try:
        logger.info(f"Processing PDF file: {file.filename} for user: {user_id}")
        pdf_content = await file.read()
        cloudinary_url, public_id = await upload_to_cloudinary(pdf_content, user_id, file.filename)
        document_id = extract_document_id_from_public_id(public_id)

        pdf_file = io.BytesIO(pdf_content)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        text = ""
        for page in pdf_reader.pages:
            page_text = page.extract_text()
            if page_text: text += page_text + "\n"
        logger.info(f"Extracted {len(text)} characters from PDF: {file.filename}")

        vector_result = await process_text_to_vectors(text, user_id, document_id)

        return {
            "content": text,
            "cloudinary_url": cloudinary_url,
            "document_id": document_id, # Return the document_id
            "filename": file.filename,
            "filetype": file.content_type,
            "vector_status": vector_result,
            "status": "success"
        }
    except Exception as e:
        logger.error(f"Error processing PDF for user {user_id}: {str(e)}")
        if not isinstance(e, HTTPException): raise HTTPException(status_code=500, detail=f"Error processing PDF: {str(e)}")
        else: raise

@router.post("/process_ppt")
async def process_ppt(user_id: str = Body(...), file: UploadFile = File(...)):
    if not file.filename.endswith(('.ppt', '.pptx')):
        raise HTTPException(status_code=400, detail="File must be a PowerPoint presentation")
    temp_path = None
    try:
        logger.info(f"Processing PowerPoint file: {file.filename} for user: {user_id}")
        file_content = await file.read()
        cloudinary_url, public_id = await upload_to_cloudinary(file_content, user_id, file.filename)
        document_id = extract_document_id_from_public_id(public_id)

        suffix = Path(file.filename).suffix
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as temp_file:
            temp_file.write(file_content)
            temp_path = temp_file.name
        logger.info(f"Saved PPT temporarily to: {temp_path}")

        try:
            prs = Presentation(temp_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"): text += shape.text + "\n"
            logger.info(f"Extracted {len(text)} characters from PPT: {file.filename}")
        except Exception as ppt_error:
            logger.error(f"Error extracting text from PowerPoint: {str(ppt_error)}")
            raise HTTPException(status_code=500, detail=f"Error reading PowerPoint content: {str(ppt_error)}")

        vector_result = await process_text_to_vectors(text, user_id, document_id)

        return {
            "content": text,
            "cloudinary_url": cloudinary_url,
            "document_id": document_id, # Return the document_id
            "filename": file.filename,
            "filetype": file.content_type,
            "vector_status": vector_result,
            "status": "success"
        }
    except Exception as e:
        logger.error(f"Error processing PowerPoint for user {user_id}: {str(e)}")
        if not isinstance(e, HTTPException): raise HTTPException(status_code=500, detail=f"Error processing PowerPoint: {str(e)}")
        else: raise
    finally:
        if temp_path and os.path.exists(temp_path):
            try: os.unlink(temp_path); logger.info(f"Successfully deleted temporary file: {temp_path}")
            except Exception as unlink_e: logger.error(f"Error deleting temporary file {temp_path}: {unlink_e}")


# --- Updated Query Endpoint ---
@router.post("/query_notes")
async def query_notes(user_id: str = Body(...), query: str = Body(...), document_id: str = Body(...)): # Expect document_id
    """Query the vector store for a specific document using Google Gemini"""
    try:
        logger.info(f"Querying notes for user: {user_id}, document_id: {document_id}. Query: '{query}'")
        # Define document-specific vector store path
        vector_store_path = f"vector_stores/{user_id}/{document_id}/vectors"

        if not os.path.exists(vector_store_path):
             logger.warning(f"Vector store not found for user: {user_id}, document_id: {document_id} at {vector_store_path}")
             raise HTTPException(status_code=404, detail="Notes for this specific document not found or not processed yet.")

        logger.info(f"Loading vector store from: {vector_store_path}")
        try:
            vectorstore = FAISS.load_local(vector_store_path, embeddings, allow_dangerous_deserialization=True)
            logger.info(f"Successfully loaded vector store for document: {document_id}")
        except Exception as faiss_load_error:
             logger.error(f"Error loading FAISS index for document {document_id}: {faiss_load_error}")
             raise HTTPException(status_code=500, detail=f"Error loading this document's index. Try re-uploading it. Original error: {str(faiss_load_error)}")

        logger.info(f"Performing similarity search for query: '{query}' on document: {document_id}")
        docs = vectorstore.similarity_search(query, k=3)

        if not docs:
             logger.info(f"No relevant documents found for query: '{query}' in document: {document_id}")
             return {"response": "Could not find relevant information in this document for your query.", "sources": [], "status": "success"}
        logger.info(f"Found {len(docs)} relevant documents for query: '{query}' in document: {document_id}")

        context = "\n\n".join([doc.page_content for doc in docs])
        prompt = f"""Answer the following question based ONLY on the provided context from the document. If the context doesn't contain the answer, say you couldn't find the information in this document.

Context:
{context}

Question: {query}

Answer:"""

        logger.info(f"Sending prompt to Gemini for document: {document_id}")
        response = llm.invoke(prompt)
        logger.info(f"Received response from Gemini for document: {document_id}")

        return {
            "response": response.content,
            "sources": [doc.page_content for doc in docs],
            "status": "success"
        }
    except HTTPException as http_exc:
         raise http_exc
    except Exception as e:
        logger.error(f"Unexpected error querying notes for document {document_id}: {str(e)}")
        if "Pickle" in str(e) or "deserialization" in str(e):
             raise HTTPException(status_code=500, detail=f"Error processing this document's index. Please try re-uploading it. Original error: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error querying notes: {str(e)}")

# --- Clear Vector Store Endpoint (remains the same, clears all for user) ---
@router.post("/clear_vector_store")
async def clear_vector_store(user_id: str = Body(...)):
    """Clear all vector stores for a specific user"""
    try:
        logger.info(f"Clearing ALL vector stores for user: {user_id}")
        vector_store_user_dir = f"vector_stores/{user_id}" # Path to the user's main folder

        if os.path.exists(vector_store_user_dir):
            shutil.rmtree(vector_store_user_dir)
            logger.info(f"All vector stores cleared successfully for user: {user_id}")
            return {"status": "success", "message": "All vector stores cleared successfully"}
        else:
            logger.info(f"No vector stores found for user: {user_id}")
            return {"status": "success", "message": "No vector stores found for this user"}
    except Exception as e:
        logger.error(f"Error clearing vector stores for user {user_id}: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error clearing vector stores: {str(e)}")

@router.post("/process_image")
async def process_image(user_id: str = Body(...), file: UploadFile = File(...)):
    """Process image files using OCR: Upload to Cloudinary, extract text, create vectors"""
    allowed_extensions = {'.png', '.jpg', '.jpeg', '.gif'}
    file_ext = os.path.splitext(file.filename)[1].lower()

    if file_ext not in allowed_extensions:
        raise HTTPException(status_code=400, detail="File must be an image (PNG, JPG, JPEG, GIF)")

    try:
        logger.info(f"Processing image file: {file.filename} for user: {user_id}")
        # Read image content
        image_content = await file.read()

        # Upload to Cloudinary first
        cloudinary_url = await upload_to_cloudinary(image_content, user_id, file.filename)

        # Perform OCR
        image = Image.open(io.BytesIO(image_content))
        text = pytesseract.image_to_string(image)
        logger.info(f"Extracted {len(text)} characters via OCR from image: {file.filename}")

        # Process text to vectors
        vector_result = await process_text_to_vectors(text, user_id, "image")

        return {
            "content": text,
            "cloudinary_url": cloudinary_url,
            "filename": file.filename,
            "filetype": file.content_type,
            "vector_status": vector_result,
            "status": "success"
        }
    except Exception as e:
        logger.error(f"Error processing image for user {user_id}: {str(e)}")
        if not isinstance(e, HTTPException):
            raise HTTPException(status_code=500, detail=f"Error processing image: {str(e)}")
        else:
            raise # Re-raise HTTPException 